import argparse
import json
import re
import subprocess
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
}

OCR_TEXT_THRESHOLD_PER_PAGE = 25


def clean_text(value):
    value = re.sub(r"[ \t\r\f\v]+", " ", value or "")
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def parse_xml(data):
    return ET.fromstring(data)


def text_from_xml(root):
    texts = []
    for node in root.findall(".//a:t", NS):
        if node.text:
            texts.append(node.text)
    return clean_text("\n".join(texts))


def natural_key(path):
    return [int(part) if part.isdigit() else part for part in re.split(r"(\d+)", path)]


def read_zip_text(zf, name):
    try:
        return zf.read(name)
    except KeyError:
        return None


def normalize_rel_target(base_dir, target):
    path = (Path(base_dir) / target).as_posix()
    parts = []
    for part in path.split("/"):
        if part == "..":
            if parts:
                parts.pop()
        elif part and part != ".":
            parts.append(part)
    return "/".join(parts)


def pptx_slide_order(zf):
    presentation = read_zip_text(zf, "ppt/presentation.xml")
    rels = read_zip_text(zf, "ppt/_rels/presentation.xml.rels")
    if not presentation or not rels:
        return sorted(
            [name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)],
            key=natural_key,
        )

    pres_root = parse_xml(presentation)
    rel_root = parse_xml(rels)
    rid_to_target = {}
    for rel in rel_root.findall(".//rel:Relationship", NS):
        rid = rel.attrib.get("Id")
        target = rel.attrib.get("Target", "")
        if rid and target.startswith("slides/"):
            rid_to_target[rid] = normalize_rel_target("ppt", target)

    order = []
    for slide_id in pres_root.findall(".//p:sldId", NS):
        rid = slide_id.attrib.get(f"{{{NS['r']}}}id")
        target = rid_to_target.get(rid)
        if target:
            order.append(target)
    return order or sorted(
        [name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)],
        key=natural_key,
    )


def slide_related_text(zf, slide_path):
    rel_path = slide_path.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
    rel_data = read_zip_text(zf, rel_path)
    notes = ""
    chart_texts = []
    image_alts = []

    if rel_data:
        root = parse_xml(rel_data)
        for rel in root.findall(".//rel:Relationship", NS):
            target = rel.attrib.get("Target", "")
            rel_type = rel.attrib.get("Type", "")
            resolved = normalize_rel_target("ppt/slides", target)
            if "notesSlide" in rel_type:
                data = read_zip_text(zf, resolved)
                if data:
                    notes = text_from_xml(parse_xml(data))
            elif "chart" in rel_type:
                data = read_zip_text(zf, resolved)
                if data:
                    chart_text = text_from_xml(parse_xml(data))
                    if chart_text:
                        chart_texts.append(chart_text)

    slide_data = read_zip_text(zf, slide_path)
    if slide_data:
        root = parse_xml(slide_data)
        for node in root.findall(".//p:cNvPr", NS):
            descr = node.attrib.get("descr") or node.attrib.get("title")
            if descr:
                image_alts.append(clean_text(descr))

    return notes, chart_texts, image_alts


def extract_pptx_ooxml(file_path):
    slides = []
    with zipfile.ZipFile(file_path) as zf:
        for index, slide_path in enumerate(pptx_slide_order(zf), start=1):
            data = read_zip_text(zf, slide_path)
            if not data:
                continue
            root = parse_xml(data)
            text = text_from_xml(root)
            notes, chart_texts, image_alts = slide_related_text(zf, slide_path)
            title = ""
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            if lines:
                title = lines[0]
            slides.append(
                {
                    "slide": index,
                    "title": title,
                    "text": text,
                    "notes": notes,
                    "charts": chart_texts,
                    "image_alt_text": image_alts,
                    "source": "ooxml",
                }
            )
    return slides


def extract_pptx_python_pptx(file_path):
    try:
        from pptx import Presentation
    except Exception as exc:
        return [], f"python-pptx unavailable: {exc}"

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        return [], f"python-pptx failed: {exc}"

    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        chunks = []
        tables = []
        title = ""
        if slide.shapes.title and getattr(slide.shapes.title, "text", None):
            title = clean_text(slide.shapes.title.text)

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(clean_text(shape.text))
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append([clean_text(cell.text) for cell in row.cells])
                tables.append(rows)

        slides.append(
            {
                "slide": index,
                "title": title,
                "text": clean_text("\n".join(chunks)),
                "tables": tables,
                "notes": "",
                "charts": [],
                "image_alt_text": [],
                "source": "python-pptx",
            }
        )
    return slides, ""


def merge_slides(primary, fallback):
    merged = []
    max_len = max(len(primary), len(fallback))
    for i in range(max_len):
        base = primary[i] if i < len(primary) else {}
        fb = fallback[i] if i < len(fallback) else {}
        slide = dict(fb)
        slide.update({k: v for k, v in base.items() if v not in (None, "", [], {})})

        for field in ("notes", "charts", "image_alt_text"):
            if not slide.get(field) and fb.get(field):
                slide[field] = fb[field]
        if not slide.get("text") and fb.get("text"):
            slide["text"] = fb["text"]
        if not slide.get("title") and fb.get("title"):
            slide["title"] = fb["title"]
        slide["slide"] = i + 1
        merged.append(slide)
    return merged


def extract_pptx(file_path):
    pptx_slides, warning = extract_pptx_python_pptx(file_path)
    ooxml_slides = extract_pptx_ooxml(file_path)
    pptx_chars = sum(len(s.get("text", "")) for s in pptx_slides)
    ooxml_chars = sum(len(s.get("text", "")) for s in ooxml_slides)

    used_fallback = bool(warning) or len(pptx_slides) != len(ooxml_slides) or pptx_chars < max(80, int(ooxml_chars * 0.6))
    slides = merge_slides(pptx_slides, ooxml_slides) if used_fallback else merge_slides(pptx_slides, ooxml_slides)
    warnings = []
    if warning:
        warnings.append(warning)
    if used_fallback:
        warnings.append("OOXML zip data was merged to improve PPTX extraction coverage.")
    return {
        "type": "pptx",
        "file": str(Path(file_path).resolve()),
        "slide_count": len(slides),
        "slides": slides,
        "warnings": warnings,
    }


def run_pdf_ocr(file_path, language, max_pages):
    script = Path(__file__).with_name("pdf-ocr.ps1")
    command = [
        "powershell.exe",
        "-NoProfile",
        "-ExecutionPolicy",
        "Bypass",
        "-File",
        str(script),
        "-PdfPath",
        str(Path(file_path).resolve()),
        "-Language",
        language or "auto",
        "-MaxPages",
        str(max_pages or 0),
    ]
    completed = subprocess.run(
        command,
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    return json.loads(completed.stdout)


def should_ocr_pdf(pages):
    if not pages:
        return False
    total_chars = sum(len(page.get("text", "")) for page in pages)
    empty_pages = sum(1 for page in pages if not page.get("text"))
    average_chars = total_chars / len(pages)
    return empty_pages / len(pages) > 0.5 or average_chars < OCR_TEXT_THRESHOLD_PER_PAGE


def merge_pdf_ocr_pages(pages, ocr_result):
    ocr_by_page = {page["page"]: page for page in ocr_result.get("pages", [])}
    merged = []
    for page in pages:
        ocr_page = ocr_by_page.get(page["page"])
        original_text = page.get("text", "")
        ocr_text = clean_text((ocr_page or {}).get("text", ""))
        source = "text"

        if ocr_text and (not original_text or len(ocr_text) > len(original_text) * 1.4):
            text = ocr_text
            source = "ocr"
        else:
            text = original_text

        merged.append(
            {
                "page": page["page"],
                "text": text,
                "source": source,
                "text_extraction": original_text,
                "ocr_text": ocr_text,
            }
        )
    return merged


def extract_pdf(file_path, ocr_scanned_pdf=False, ocr_language="auto", max_ocr_pages=80):
    warnings = []
    pages = []
    try:
        from pypdf import PdfReader
    except Exception as exc:
        if ocr_scanned_pdf:
            try:
                ocr_result = run_pdf_ocr(file_path, ocr_language, max_ocr_pages)
                pages = [
                    {
                        "page": page["page"],
                        "text": clean_text(page.get("text", "")),
                        "source": "ocr",
                        "text_extraction": "",
                        "ocr_text": clean_text(page.get("text", "")),
                    }
                    for page in ocr_result.get("pages", [])
                ]
                return {
                    "type": "pdf",
                    "file": str(Path(file_path).resolve()),
                    "page_count": ocr_result.get("pageCount", len(pages)),
                    "pages": pages,
                    "ocr": {
                        "attempted": True,
                        "language": ocr_result.get("language"),
                        "processed_pages": ocr_result.get("processedPages", 0),
                    },
                    "warnings": [
                        f"pypdf unavailable in this Python runtime: {exc}",
                        "Fell back to local Windows PDF rendering OCR.",
                    ],
                }
            except Exception as ocr_exc:
                warning = (
                    f"pypdf unavailable in this Python runtime: {exc}; "
                    f"scanned PDF OCR also failed: {ocr_exc}"
                )
        else:
            warning = f"pypdf unavailable in this Python runtime: {exc}"
        return {
            "type": "pdf",
            "file": str(Path(file_path).resolve()),
            "page_count": 0,
            "pages": [],
            "ocr": {"attempted": False, "language": None, "processed_pages": 0},
            "warnings": [warning],
        }

    reader = PdfReader(file_path)
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = clean_text(page.extract_text() or "")
        except Exception as exc:
            text = ""
            warnings.append(f"Page {index} text extraction failed: {exc}")
        pages.append({"page": index, "text": text, "source": "text"})

    empty_pages = sum(1 for p in pages if not p["text"])
    if pages and empty_pages / len(pages) > 0.5:
        warnings.append("Many pages contain little or no text. This may be a scanned PDF.")

    ocr_result = None
    if ocr_scanned_pdf and should_ocr_pdf(pages):
        try:
            ocr_result = run_pdf_ocr(file_path, ocr_language, max_ocr_pages)
            pages = merge_pdf_ocr_pages(pages, ocr_result)
            warnings.append(
                f"Scanned PDF OCR ran locally with language {ocr_result.get('language', ocr_language)} "
                f"on {ocr_result.get('processedPages', 0)} page(s)."
            )
            if ocr_result.get("processedPages", 0) < len(reader.pages):
                warnings.append("OCR page limit was reached; some pages were not OCR processed.")
        except Exception as exc:
            warnings.append(f"Scanned PDF OCR failed: {exc}")

    return {
        "type": "pdf",
        "file": str(Path(file_path).resolve()),
        "page_count": len(pages),
        "pages": pages,
        "ocr": {
            "attempted": bool(ocr_result),
            "language": (ocr_result or {}).get("language"),
            "processed_pages": (ocr_result or {}).get("processedPages", 0),
        },
        "warnings": warnings,
    }


def result_to_markdown(result, max_chars):
    path = Path(result["file"])
    lines = [
        f"# Courseware Extraction: {path.name}",
        "",
        f"- Type: {result['type']}",
    ]

    if result.get("slide_count") is not None:
        lines.append(f"- Slides: {result['slide_count']}")
    if result.get("page_count") is not None:
        lines.append(f"- Pages: {result['page_count']}")
    for warning in result.get("warnings", []):
        lines.append(f"- Warning: {warning}")
    lines.append("")

    if result["type"] == "pptx":
        for slide in result.get("slides", []):
            title = slide.get("title") or f"Slide {slide['slide']}"
            lines.append(f"## Slide {slide['slide']}: {title}")
            if slide.get("text"):
                lines.append(slide["text"])
            if slide.get("notes"):
                lines.append("")
                lines.append(f"Speaker notes: {slide['notes']}")
            if slide.get("tables"):
                lines.append("")
                lines.append("Tables:")
                for table in slide["tables"]:
                    for row in table:
                        lines.append("- " + " | ".join(cell or "" for cell in row))
            if slide.get("charts"):
                lines.append("")
                lines.append("Charts:")
                for chart in slide["charts"]:
                    lines.append(f"- {chart}")
            if slide.get("image_alt_text"):
                lines.append("")
                lines.append("Image alt text:")
                for alt in slide["image_alt_text"]:
                    lines.append(f"- {alt}")
            lines.append("")
    else:
        for page in result.get("pages", []):
            source = page.get("source", "text")
            lines.append(f"## Page {page['page']} ({source})")
            lines.append(page.get("text") or "[No extractable text]")
            lines.append("")

    markdown = "\n".join(lines).strip()
    truncated = len(markdown) > max_chars
    if truncated:
        markdown = markdown[:max_chars].rstrip() + "\n\n[Truncated due to max_chars limit.]"
    result["markdown_truncated"] = truncated
    result["markdown"] = markdown
    return result


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", required=True)
    parser.add_argument("--max-chars", type=int, default=60000)
    parser.add_argument("--ocr-scanned-pdf", action="store_true")
    parser.add_argument("--ocr-language", default="auto")
    parser.add_argument("--max-ocr-pages", type=int, default=80)
    args = parser.parse_args()

    file_path = Path(args.file)
    if not file_path.exists():
        raise FileNotFoundError(str(file_path))

    ext = file_path.suffix.lower()
    if ext in (".pptx", ".ppt"):
        result = extract_pptx(file_path)
    elif ext == ".pdf":
        result = extract_pdf(
            file_path,
            ocr_scanned_pdf=args.ocr_scanned_pdf,
            ocr_language=args.ocr_language,
            max_ocr_pages=args.max_ocr_pages,
        )
    else:
        raise ValueError(f"Unsupported courseware type: {ext}")

    print(json.dumps(result_to_markdown(result, args.max_chars), ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        sys.exit(1)
